import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class PdfToPptxConverter:
    def __init__(self, output_filename="presentation.pptx", ollama_processor=None, theme="default"):
        self.prs = Presentation()
        self.output_filename = output_filename
        self.ollama_processor = ollama_processor
        self.theme = theme

        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)

        self.setup_theme_colors()

    def setup_theme_colors(self):
        if self.theme == "corporate":
            self.title_color = RGBColor(18, 52, 86)
            self.accent_color = RGBColor(64, 119, 176)
            self.text_color = RGBColor(50, 50, 50)
            self.background_color = RGBColor(242, 242, 242)
            self.title_font_size = Pt(36)
            self.subtitle_font_size = Pt(20)
            self.header_font_size = Pt(28)
            self.content_font_size = Pt(16)
        elif self.theme == "minimal":
            self.title_color = RGBColor(0, 0, 0)
            self.accent_color = RGBColor(204, 0, 0)
            self.text_color = RGBColor(40, 40, 40)
            self.background_color = RGBColor(255, 255, 255)
            self.title_font_size = Pt(38)
            self.subtitle_font_size = Pt(22)
            self.header_font_size = Pt(30)
            self.content_font_size = Pt(18)
        else:
            self.title_color = RGBColor(44, 86, 151)
            self.accent_color = RGBColor(0, 129, 198)
            self.text_color = RGBColor(68, 68, 68)
            self.background_color = RGBColor(250, 250, 250)
            self.title_font_size = Pt(36)
            self.subtitle_font_size = Pt(22)
            self.header_font_size = Pt(28)
            self.content_font_size = Pt(18)

    def _add_title_slide(self, title, subtitle=None):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        title_shape.left = Inches(0.5)
        title_shape.top = Inches(2.5)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(2.0)

        text_frame = title_shape.text_frame
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.font.size = self.title_font_size
        p.font.bold = True
        p.font.color.rgb = self.title_color
        p.alignment = PP_ALIGN.CENTER

        if subtitle:
            if len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle

                subtitle_shape.left = Inches(0.5)
                subtitle_shape.top = Inches(4.8)
                subtitle_shape.width = Inches(12.33)
                subtitle_shape.height = Inches(1.0)

                subtitle_frame = subtitle_shape.text_frame
                subtitle_frame.margin_left = Inches(0.2)
                subtitle_frame.margin_right = Inches(0.2)
                subtitle_frame.margin_top = Inches(0.1)
                subtitle_frame.margin_bottom = Inches(0.1)

                p = subtitle_frame.paragraphs[0]
                p.font.size = self.subtitle_font_size
                p.font.color.rgb = self.accent_color
                p.alignment = PP_ALIGN.CENTER

        if self.theme == "minimal":
            left = Inches(3.5)
            top = Inches(4.5)
            width = Inches(6.33)
            height = Inches(0.05)
            line = slide.shapes.add_shape(1, left, top, width, height)
            line.fill.solid()
            line.fill.fore_color.rgb = self.accent_color
            line.line.fill.background()

        return slide

    def _add_section_slide(self, title):
        slide_layout = self.prs.slide_layouts[2] if len(self.prs.slide_layouts) > 2 else self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        title_shape.left = Inches(0.5)
        title_shape.top = Inches(1.0)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(1.5)

        text_frame = title_shape.text_frame
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)

        p = text_frame.paragraphs[0]
        p.font.size = self.header_font_size
        p.font.bold = True
        p.font.color.rgb = self.title_color
        p.alignment = PP_ALIGN.CENTER

        if self.theme == "corporate":
            left = Inches(0)
            top = Inches(0)
            width = Inches(1.0)
            height = Inches(7.5)
            sidebar = slide.shapes.add_shape(1, left, top, width, height)
            sidebar.fill.solid()
            sidebar.fill.fore_color.rgb = self.accent_color
            sidebar.line.fill.background()
        elif self.theme == "minimal":
            left = Inches(0.5)
            top = Inches(2.8)
            width = Inches(0.1)
            height = Inches(2.5)
            line = slide.shapes.add_shape(1, left, top, width, height)
            line.fill.solid()
            line.fill.fore_color.rgb = self.accent_color
            line.line.fill.background()

        return slide

    def _add_content_slide(self, title, content_points):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.3)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(1.2)

        title_frame = title_shape.text_frame
        title_frame.margin_left = Inches(0.2)
        title_frame.margin_right = Inches(0.2)
        title_frame.margin_top = Inches(0.1)
        title_frame.margin_bottom = Inches(0.1)

        p = title_frame.paragraphs[0]
        p.font.size = self.header_font_size
        p.font.bold = True
        p.font.color.rgb = self.title_color

        if len(slide.placeholders) > 1:
            content = slide.placeholders[1]
        else:
            left = Inches(0.8)
            top = Inches(1.8)
            width = Inches(11.73)
            height = Inches(5.2)
            content = slide.shapes.add_textbox(left, top, width, height)

        content.left = Inches(0.8)
        content.top = Inches(1.8)
        content.width = Inches(11.73)
        content.height = Inches(5.2)

        text_frame = content.text_frame
        text_frame.clear()

        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_bottom = Inches(0.2)
        text_frame.word_wrap = True

        for i, point in enumerate(content_points):
            if point.strip():
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = f"• {point.strip()}"
                p.font.size = self.content_font_size
                p.font.color.rgb = self.text_color
                p.level = 0
                p.space_after = Pt(6)

        if self.theme == "corporate":
            left = Inches(0)
            top = Inches(7.0)
            width = Inches(13.33)
            height = Inches(0.5)
            footer = slide.shapes.add_shape(1, left, top, width, height)
            footer.fill.solid()
            footer.fill.fore_color.rgb = self.accent_color
            footer.line.fill.background()
        elif self.theme == "minimal":
            left = Inches(0.8)
            top = Inches(1.6)
            width = Inches(11.73)
            height = Inches(0.03)
            line = slide.shapes.add_shape(1, left, top, width, height)
            line.fill.solid()
            line.fill.fore_color.rgb = self.accent_color
            line.line.fill.background()

        return slide

    def _add_content_slide_with_image(self, title, content_points, image_path):
        """
        Adds a slide with text and image side by side with automatic sizing.
        """
        slide_layout = self.prs.slide_layouts[1]  # Layout with title and content
        slide = self.prs.slides.add_slide(slide_layout)

        # Configure title
        title_shape = slide.shapes.title
        title_shape.text = title

        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.3)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(1.0)

        p = title_shape.text_frame.paragraphs[0]
        p.font.size = self.header_font_size
        p.font.bold = True
        p.font.color.rgb = self.title_color

        # If there's no image or the image doesn't exist, create normal content slide
        if not image_path or not os.path.exists(image_path):
            print(f"Image not found: {image_path}")
            return self._add_content_slide(title, content_points)

        # Check image dimensions
        try:
            from PIL import Image
            img = Image.open(image_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height
            print(f"Adding image: {image_path}, dimensions: {img_width}x{img_height}")
        except Exception as e:
            print(f"Error analyzing image: {e}")
            return self._add_content_slide(title, content_points)

        # Configure slide division: 60% text, 40% image
        text_left = Inches(0.8)
        text_top = Inches(1.5)
        text_width = Inches(6.5)  # Reduced to give more space to the image
        text_height = Inches(5.0)

        # Add text area
        content = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        text_frame = content.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)

        # Add content with bullet points
        for i, point in enumerate(content_points):
            if point and isinstance(point, str):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = point.strip()
                p.font.size = self.content_font_size
                p.font.color.rgb = self.text_color
                p.level = 0  # Bullet level
                p.space_after = Pt(6)  # Spacing after paragraph

        # Calculate ideal dimensions for the image
        img_left = Inches(7.5)  # Positioned further left to give adequate space
        img_top = Inches(1.5)
        img_max_width = Inches(5.0)  # Increased to allow larger images
        img_max_height = Inches(5.0)

        # Calculate dimensions preserving aspect ratio
        if aspect_ratio > 1:  # Image is wider than tall
            img_width = min(img_max_width, Inches(5.0))
            img_height = img_width / aspect_ratio
            if img_height > img_max_height:
                img_height = img_max_height
                img_width = img_height * aspect_ratio
        else:  # Image is taller than wide
            img_height = min(img_max_height, Inches(5.0))
            img_width = img_height * aspect_ratio
            if img_width > img_max_width:
                img_width = img_max_width
                img_height = img_width / aspect_ratio

        # Center the image in the designated area
        img_left = img_left + (img_max_width - img_width) / 2
        img_top = img_top + (img_max_height - img_height) / 2

        # Add the image to the slide
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width, height=img_height)
            print(f"Image added successfully: {image_path}")
        except Exception as e:
            print(f"Error adding image to slide: {e}")

        return slide

    def _add_table_slide(self, title, table_data):
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.3)
        title_shape.width = Inches(12.33)
        title_shape.height = Inches(1.0)

        p = title_shape.text_frame.paragraphs[0]
        p.font.size = self.header_font_size
        p.font.bold = True
        p.font.color.rgb = self.title_color

        if table_data and len(table_data) > 0:
            rows = len(table_data)
            cols = max(len(row) for row in table_data) if table_data else 1

            left = Inches(0.8)
            top = Inches(1.8)
            width = Inches(11.73)
            height = Inches(4.5)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            for row_idx, row_data in enumerate(table_data):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < cols:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(cell_data)

                        p = cell.text_frame.paragraphs[0]
                        p.font.size = Pt(14)
                        p.font.color.rgb = self.text_color

                        if row_idx == 0:
                            p.font.bold = True
                            p.font.color.rgb = self.title_color

        return slide

    def _detect_tables(self, text_content):
        lines = text_content.strip().split('\n')
        table_candidate_lines = 0

        for line in lines:
            if re.search(r'\s{3,}|\t', line) or line.count('|') >= 2:
                table_candidate_lines += 1

        return table_candidate_lines >= 3

    def _process_table_data(self, content_lines):
        if not content_lines:
            return []

        table_data = []

        for line in content_lines:
            if re.search(r'\s{3,}|\t', line):
                cells = re.split(r'\s{3,}|\t', line)
                cells = [cell.strip() for cell in cells if cell.strip()]

                if cells:
                    table_data.append(cells)
            elif line.count('|') >= 2:
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                if cells:
                    table_data.append(cells)

        return table_data

    def create_presentation(self, document_structure, image_data=None):
        if not isinstance(document_structure, dict):
            document_structure = self._convert_to_structure(document_structure)

        # Initial validation of image_data
        if image_data and not isinstance(image_data, list):
            print(f"Warning: invalid image_data, expected format: list, received: {type(image_data)}")
            image_data = []

        title = document_structure.get('title', 'Document')
        subtitle = document_structure.get('subtitle', '')

        if document_structure.get('version'):
            if subtitle:
                subtitle += f" | Version: {document_structure['version']}"
            else:
                subtitle = f"Version: {document_structure['version']}"

        if document_structure.get('date'):
            if subtitle:
                subtitle += f" | {document_structure['date']}"
            else:
                subtitle = document_structure['date']

        # Create title slide
        self._add_title_slide(title, subtitle)

        # Process sections
        for section in document_structure.get('sections', []):
            section_title = section.get('title', '')
            content = section.get('content', [])

            # Check if there are images associated with this section
            section_image = None
            if image_data and section.get('has_images'):
                img_info = section.get('image_info', {})
                relevant_images = img_info.get('relevant_images', [])

                if relevant_images and len(relevant_images) > 0:
                    # Get index of the first relevant image
                    img_idx = relevant_images[0]
                    if 0 <= img_idx < len(image_data):
                        img_path = image_data[img_idx].get('path') if isinstance(image_data[img_idx], dict) else \
                        image_data[img_idx]
                        if os.path.exists(img_path):
                            section_image = img_path
                            print(f"Associating image {img_path} with section '{section_title}'")

            # For sections without explicitly associated images, search for images by page correspondence
            if not section_image and image_data:
                for img in image_data:
                    if isinstance(img, dict) and 'path' in img:
                        section_image = img['path']
                        print(f"Associating default image {section_image} with section '{section_title}'")
                        break

            # Add slide according to content type
            if content and isinstance(content, list) and len(content) > 0:
                # Check if the content is a table
                if len(content) == 1 and self._detect_tables(content[0]):
                    table_data = self._process_table_data([content[0]])
                    self._add_table_slide(section_title, table_data)
                else:
                    # Add slide with or without image
                    if section_image:
                        self._add_content_slide_with_image(section_title, content, section_image)
                    else:
                        self._add_content_slide(section_title, content)

        # Save the presentation
        self.prs.save(self.output_filename)
        return self.output_filename

    def _is_image_relevant(self, section_content, image_path):
        """
        Checks if the image is relevant to the section content.
        """
        keywords = ["diagram", "graphic", "figure", "image", "visualization"]
        content_text = " ".join(section_content).lower()

        # Check if the content mentions keywords related to images
        if any(keyword in content_text for keyword in keywords):
            return True

        # Add other checks, such as similarity analysis, if necessary
        return False

    def _convert_to_structure(self, text_content):
        import re

        if isinstance(text_content, dict):
            return text_content

        structure = {
            "title": "",
            "subtitle": "",
            "sections": []
        }

        lines = text_content.strip().split('\n')
        current_section = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            if re.match(r'^#+\s', line) or (line.isupper() and len(line) < 100):
                if not structure["title"]:
                    structure["title"] = line.lstrip('#').strip()
                    continue

                section_title = line.lstrip('#').strip()
                if not section_title:
                    continue

                current_section = {
                    "title": section_title,
                    "content": []
                }
                structure["sections"].append(current_section)
            elif current_section and (line.startswith('*') or line.startswith('-') or re.match(r'^\d+\.', line)):
                item_text = line.lstrip('*-0123456789. \t')
                current_section["content"].append(item_text)

        if not structure["sections"] and text_content:
            structure["sections"] = [{
                "title": "General Information",
                "content": [line.strip() for line in lines if line.strip()]
            }]

        return structure

    def _select_image_for_section(self, section, image_data):
        """
        Selects the most suitable image for a section based on analysis.

        Args:
            section (dict): The document section
            image_data (list): List of available image data

        Returns:
            dict or None: Selected image data or None if no suitable image
        """
        if not image_data or not section.get("has_images", False):
            return None

        image_info = section.get("image_info", {})
        relevant_images = image_info.get("relevant_images", [])

        if not relevant_images:
            # Alternative strategy: look for image references in text
            content_text = " ".join(section.get("content", [])).lower()
            if any(keyword in content_text for keyword in ["figure", "image", "graphic", "diagram", "illustration"]):
                # If there are image references in the text, select the first available image
                if image_data:
                    return image_data[0]
            return None

        # Select the first relevant image that exists in image_data
        for img_index in relevant_images:
            if 0 <= img_index < len(image_data):
                return image_data[img_index]

        return None
